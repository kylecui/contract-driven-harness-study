# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-pptx>=1.0.0,<2",
# ]
# ///
"""Build a PPTX deck from a structured JSON deck spec.

Usage:
  uv run scripts/build_deck.py deck_spec.json --out output/deck.pptx
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WIDE_W = 13.333
WIDE_H = 7.5


def hex_to_rgb(value: str) -> RGBColor:
    value = (value or "000000").strip().lstrip("#")
    if not re.fullmatch(r"[0-9a-fA-F]{6}", value):
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def theme(spec: dict[str, Any]) -> dict[str, Any]:
    t = spec.get("theme") or {}
    return {
        "primary": t.get("primary", "1E2761"),
        "secondary": t.get("secondary", "CADCFC"),
        "accent": t.get("accent", "F96167"),
        "background": t.get("background", "FFFFFF"),
        "text": t.get("text", "1F2937"),
        "font_face": t.get("font_face", "Microsoft YaHei"),
    }


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size=18,
    bold=False,
    color="1F2937",
    font_face="Microsoft YaHei",
    align=PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.name = font_face
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return box


def add_paragraphs(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size=18,
    color="1F2937",
    font_face="Microsoft YaHei",
    bullet=True,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        if bullet:
            # python-pptx does not expose every bullet property cleanly; native paragraphs are enough for default viewers.
            p.text = f"• {p.text}"
        p.font.name = font_face
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(color)
        p.space_after = Pt(8)
    return box


def add_footer(slide, index: int, th: dict[str, Any]) -> None:
    add_textbox(
        slide,
        f"{index}",
        12.4,
        7.05,
        0.5,
        0.2,
        font_size=9,
        color="6B7280",
        font_face=th["font_face"],
        align=PP_ALIGN.RIGHT,
    )


def add_top_rule(slide, th: dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(WIDE_W),
        Inches(0.12),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(th["primary"])
    shape.line.fill.background()


def title_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["primary"])
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(6.9),
        Inches(WIDE_W),
        Inches(0.6),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(th["accent"])
    accent.line.fill.background()
    add_textbox(
        slide,
        spec.get("title", "Untitled"),
        0.85,
        1.55,
        11.5,
        1.3,
        font_size=38,
        bold=True,
        color="FFFFFF",
        font_face=th["font_face"],
    )
    if spec.get("subtitle"):
        add_textbox(
            slide,
            spec["subtitle"],
            0.9,
            3.05,
            10.8,
            0.8,
            font_size=20,
            color=th["secondary"],
            font_face=th["font_face"],
        )
    add_footer(slide, idx, {**th, "font_face": th["font_face"]})


def section_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["primary"])
    add_textbox(
        slide,
        f"{idx:02d}",
        0.85,
        0.7,
        1.0,
        0.5,
        font_size=16,
        bold=True,
        color=th["secondary"],
        font_face=th["font_face"],
    )
    add_textbox(
        slide,
        spec.get("title", "Section"),
        0.85,
        2.3,
        11.5,
        0.9,
        font_size=34,
        bold=True,
        color="FFFFFF",
        font_face=th["font_face"],
    )
    if spec.get("subtitle"):
        add_textbox(
            slide,
            spec["subtitle"],
            0.9,
            3.45,
            10.8,
            0.6,
            font_size=18,
            color=th["secondary"],
            font_face=th["font_face"],
        )
    add_footer(slide, idx, th)


def bullets_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["background"])
    add_top_rule(slide, th)
    add_textbox(
        slide,
        spec.get("title", "Untitled"),
        0.65,
        0.55,
        12.0,
        0.55,
        font_size=28,
        bold=True,
        color=th["primary"],
        font_face=th["font_face"],
    )
    y = 1.35
    if spec.get("body"):
        add_textbox(
            slide,
            spec["body"],
            0.7,
            y,
            11.5,
            0.55,
            font_size=15,
            color="4B5563",
            font_face=th["font_face"],
        )
        y += 0.75
    items = [str(x) for x in spec.get("items", [])][:6]
    # cards
    for i, item in enumerate(items):
        yy = y + i * 0.72
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            Inches(yy),
            Inches(11.8),
            Inches(0.52),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb("F8FAFC")
        card.line.color.rgb = hex_to_rgb(th["secondary"])
        add_textbox(
            slide,
            item,
            1.05,
            yy + 0.09,
            11.0,
            0.32,
            font_size=16,
            color=th["text"],
            font_face=th["font_face"],
        )
    add_footer(slide, idx, th)


def two_column_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["background"])
    add_top_rule(slide, th)
    add_textbox(
        slide,
        spec.get("title", "Untitled"),
        0.65,
        0.55,
        12.0,
        0.55,
        font_size=28,
        bold=True,
        color=th["primary"],
        font_face=th["font_face"],
    )
    columns = [
        (
            0.75,
            spec.get("left_title", "Left"),
            spec.get("left_items", []),
            th["secondary"],
        ),
        (
            6.8,
            spec.get("right_title", "Right"),
            spec.get("right_items", []),
            th["accent"],
        ),
    ]
    for x, title, items, color in columns:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.45),
            Inches(5.75),
            Inches(5.15),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb("F8FAFC")
        card.line.color.rgb = hex_to_rgb(color)
        add_textbox(
            slide,
            title,
            x + 0.35,
            1.75,
            5.1,
            0.38,
            font_size=20,
            bold=True,
            color=th["primary"],
            font_face=th["font_face"],
        )
        add_paragraphs(
            slide,
            [str(i) for i in items][:6],
            x + 0.35,
            2.35,
            5.05,
            3.8,
            font_size=15,
            color=th["text"],
            font_face=th["font_face"],
        )
    add_footer(slide, idx, th)


def comparison_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["background"])
    add_top_rule(slide, th)
    add_textbox(
        slide,
        spec.get("title", "Comparison"),
        0.65,
        0.55,
        12.0,
        0.55,
        font_size=28,
        bold=True,
        color=th["primary"],
        font_face=th["font_face"],
    )
    cols = spec.get("columns", [])[:3]
    if not cols:
        cols = [{"title": "Option A", "items": []}, {"title": "Option B", "items": []}]
    gap = 0.25
    card_w = (11.9 - gap * (len(cols) - 1)) / len(cols)
    for i, col in enumerate(cols):
        x = 0.75 + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.45),
            Inches(card_w),
            Inches(5.15),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb("F8FAFC")
        card.line.color.rgb = hex_to_rgb(
            th["secondary"] if i % 2 == 0 else th["accent"]
        )
        add_textbox(
            slide,
            str(col.get("title", f"Column {i + 1}")),
            x + 0.25,
            1.75,
            card_w - 0.5,
            0.4,
            font_size=18,
            bold=True,
            color=th["primary"],
            font_face=th["font_face"],
            align=PP_ALIGN.CENTER,
        )
        add_paragraphs(
            slide,
            [str(x) for x in col.get("items", [])][:6],
            x + 0.25,
            2.35,
            card_w - 0.5,
            3.8,
            font_size=14,
            color=th["text"],
            font_face=th["font_face"],
        )
    add_footer(slide, idx, th)


def process_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["background"])
    add_top_rule(slide, th)
    add_textbox(
        slide,
        spec.get("title", "Process"),
        0.65,
        0.55,
        12.0,
        0.55,
        font_size=28,
        bold=True,
        color=th["primary"],
        font_face=th["font_face"],
    )
    steps = spec.get("steps", [])[:5]
    if not steps:
        steps = [{"title": "Step 1", "body": ""}]
    gap = 0.25
    card_w = (11.9 - gap * (len(steps) - 1)) / len(steps)
    for i, step in enumerate(steps):
        x = 0.75 + i * (card_w + gap)
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + 0.15),
            Inches(1.65),
            Inches(0.6),
            Inches(0.6),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(th["accent"])
        circle.line.fill.background()
        add_textbox(
            slide,
            str(i + 1),
            x + 0.15,
            1.78,
            0.6,
            0.3,
            font_size=16,
            bold=True,
            color="FFFFFF",
            font_face=th["font_face"],
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            str(step.get("title", f"Step {i + 1}")),
            x,
            2.45,
            card_w,
            0.5,
            font_size=17,
            bold=True,
            color=th["primary"],
            font_face=th["font_face"],
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            str(step.get("body", "")),
            x + 0.05,
            3.1,
            card_w - 0.1,
            1.3,
            font_size=13,
            color=th["text"],
            font_face=th["font_face"],
            align=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            add_textbox(
                slide,
                "→",
                x + card_w - 0.05,
                1.72,
                0.4,
                0.3,
                font_size=22,
                bold=True,
                color=th["secondary"],
                font_face=th["font_face"],
                align=PP_ALIGN.CENTER,
            )
    add_footer(slide, idx, th)


def quote_slide(prs, spec: dict[str, Any], th: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["primary"])
    add_textbox(
        slide,
        spec.get("title", "Key Message"),
        0.85,
        0.65,
        11.8,
        0.5,
        font_size=20,
        bold=True,
        color=th["secondary"],
        font_face=th["font_face"],
    )
    quote = spec.get("quote", "")
    add_textbox(
        slide,
        f"“{quote}”",
        1.15,
        2.05,
        11.0,
        1.6,
        font_size=30,
        bold=True,
        color="FFFFFF",
        font_face=th["font_face"],
        align=PP_ALIGN.CENTER,
    )
    if spec.get("attribution"):
        add_textbox(
            slide,
            f"— {spec['attribution']}",
            1.15,
            4.2,
            11.0,
            0.4,
            font_size=15,
            color=th["secondary"],
            font_face=th["font_face"],
            align=PP_ALIGN.CENTER,
        )
    add_footer(slide, idx, th)


def build(spec: dict[str, Any], out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)
    th = theme(spec)
    meta = spec.get("meta") or {}
    prs.core_properties.title = meta.get("title", "Generated deck")
    prs.core_properties.author = meta.get("author", "ppt-writer")
    slides = spec.get("slides") or []
    if not slides:
        slides = [
            {
                "type": "title",
                "title": meta.get("title", "Generated deck"),
                "subtitle": meta.get("subtitle", ""),
            }
        ]

    dispatch = {
        "title": title_slide,
        "section": section_slide,
        "bullets": bullets_slide,
        "two-column": two_column_slide,
        "comparison": comparison_slide,
        "process": process_slide,
        "quote": quote_slide,
    }
    for idx, slide_spec in enumerate(slides, start=1):
        typ = slide_spec.get("type") or slide_spec.get("layout", "bullets")
        fn = dispatch.get(typ, bullets_slide)
        fn(prs, slide_spec, th, idx)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Build a PPTX deck from JSON deck spec."
    )
    parser.add_argument("spec", type=Path, help="Deck spec JSON")
    parser.add_argument("--out", type=Path, required=True, help="Output .pptx path")
    args = parser.parse_args()
    try:
        spec = json.loads(args.spec.read_text(encoding="utf-8"))
        build(spec, args.out)
        print(
            json.dumps(
                {"output": str(args.out), "slides": len(spec.get("slides", []))},
                ensure_ascii=False,
                indent=2,
            )
        )
        return 0
    except Exception as exc:  # noqa: BLE001
        print(f"Error: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
